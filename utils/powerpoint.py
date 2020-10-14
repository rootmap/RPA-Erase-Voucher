from pptx import Presentation
from pptx.util import Inches, Pt

from utils.pd2ppt.pd2ppt import DataFrame2PPT


class Powerpoint(DataFrame2PPT):

    def create_powerpoint(self, file_name):
        ppt = Presentation()
        self.save_ppt(ppt, file_name)

    def save_ppt(self, ppt, file_name):
        ppt.save(file_name)

    def open_powerpoint(self, file_name):
        ppt = Presentation(file_name)
        return ppt

    def create_slide(self, ppt, layout):
        return ppt.slides.add_slide(layout)

    def add_blank_layout(self, ppt):
        ppt_layout = ppt.slide_layouts[6]
        current_slide = self.create_slide(ppt, ppt_layout)
        return current_slide

    def add_text(self, slide, text, bold, font_size, pos_left, pos_right):
        area_left = Inches(pos_left)
        area_right = Inches(pos_right)
        area_height_width = Inches(1)
        text_area = slide.shapes.add_textbox(area_left, area_right, area_height_width, area_height_width)
        text_frame = text_area.text_frame

        current_paragraph = text_frame.add_paragraph()
        current_paragraph.text = text

        if bold == True:
            current_paragraph.font.bold = True

        current_paragraph.font.size = Pt(font_size)

    def add_image(self, slide, image, pos_left, pos_top, width, *height):
        area_left = Inches(pos_left)
        area_top = Inches(pos_top)
        area_width = Inches(width)

        if len(height) > 0:
            area_height = Inches(height[0])
            slide.shapes.add_picture(image, area_left, area_top, area_width, area_height)
        else:
            slide.shapes.add_picture(image, area_left, area_top, area_width)

    def add_table(self, slide, df, pos_left=None, pos_top=None, width=None, *height):
        if None in {pos_left, pos_top, width}:
            pos_left = 1
            pos_top = 3
            width = 5
        if len(height) > 0:
            height = height[0]
        else:
            height = width

        super().df_to_table(slide, df, pos_left, pos_top, width, height)


    def add_df_to_new_slide(self, filename, df, **kwargs):
        if kwargs is not None:
            super().df_to_powerpoint(filename, df, **kwargs)
        else:
            super().df_to_powerpoint(filename, df, left=1, top=3, width=5, height=5)
